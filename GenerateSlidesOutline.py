import os
import re
import base64
import argparse
from pathlib import Path
import pptx
from pptx.util import Inches
from openai import OpenAI
import logging
import dotenv
from typing import List, Dict, Any, Optional

# Load environment variables
dotenv.load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Import necessary functions from MinerU.py
from MinerU import extract_pdf_to_markdown

def parse_arguments():
    """Parse command line arguments"""
    parser = argparse.ArgumentParser(description='Generate PowerPoint slides from PDF papers')
    parser.add_argument('--model', type=str, default='openai', 
                        choices=['openai', 'gemini', 'claude', 'grok'],
                        help='LLM model to use (default: openai)')
    parser.add_argument('--pdf', type=str, default='data/Example.pdf',
                        help='Path to PDF file (default: data/Example.pdf)')
    parser.add_argument('--output', type=str, default='Generated_Slides.pptx',
                        help='Output PowerPoint file name (default: Generated_Slides.pptx)')
    parser.add_argument('--api-key', type=str, 
                        help='API key (default: read from environment variables)')
    return parser.parse_args()

def setup_api_keys(args):
    """Set up API keys"""
    # If API key is provided via command line, use it
    if args.api_key:
        os.environ["OPENAI_API_KEY"] = args.api_key
    # Otherwise check environment variables
    elif not os.environ.get("OPENAI_API_KEY"):
        # Loading from .env file is already done
        if not os.environ.get("OPENAI_API_KEY"):
            raise ValueError("Missing OpenAI API key. Please provide it via --api-key parameter or set OPENAI_API_KEY environment variable")
    
    return os.environ.get("OPENAI_API_KEY")

def get_markdown_file_path(pdf_path: str, output_dir: str) -> str:
    """
    Get the corresponding Markdown file path from PDF path
    
    Args:
        pdf_path: PDF file path
        output_dir: Output directory
        
    Returns:
        Markdown file path
    """
    pdf_base_name = Path(pdf_path).stem
    return str(Path(output_dir) / pdf_base_name / "auto" / f"{pdf_base_name}.md")

def extract_image_paths(md_text: str) -> List[str]:
    """
    Extract image paths from Markdown text
    
    Args:
        md_text: Markdown text
        
    Returns:
        List of image paths
    """
    return re.findall(r'!\[.*?\]\((.*?)\)', md_text)

def encode_image(image_path: str) -> str:
    """
    Encode image to base64
    
    Args:
        image_path: Path to the image
        
    Returns:
        Base64-encoded image
    """
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        logger.error(f"Error encoding image {image_path}: {e}")
        return ""

def prepare_image_messages(image_paths: List[str], pdf_base_name: str) -> List[Dict[str, Any]]:
    """
    Prepare messages with base64-encoded images
    
    Args:
        image_paths: List of image paths
        pdf_base_name: PDF file base name
        
    Returns:
        List of messages with images
    """
    image_messages = []
    for image_path in image_paths:
        full_image_path = str(Path("output") / pdf_base_name / "auto" / image_path)
        base64_image = encode_image(full_image_path)
        if base64_image:
            image_messages.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{base64_image}",
                }
            })
    return image_messages

def load_prompt_template(template_path: str) -> str:
    """
    Load prompt template
    
    Args:
        template_path: Path to template file
        
    Returns:
        Prompt template text
    """
    try:
        with open(template_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error loading prompt template {template_path}: {e}")
        raise

def get_openai_client(model_name: str, api_key: str) -> Any:
    """
    Get LLM client
    
    Args:
        model_name: Model name
        api_key: API key
        
    Returns:
        LLM client
    """
    return OpenAI(api_key=api_key)

def get_gemini_client(model_name: str, api_key: str) -> Any:
    """
    Get Gemini client
    """
    from google.generativeai import GenerativeAI
    return GenerativeAI.Client(api_key=os.environ.get("GEMINI_API_KEY"))

def get_grok_client(model_name: str, api_key: str) -> Any:
    """
    Get Grok client
    """
    return OpenAI(api_key=os.environ.get("GROK_API_KEY"), base_url=os.environ.get("GROK_BASE_URL"))

def generate_outline(client: Any, md_text: str, image_messages: List[Dict[str, Any]], model_name: str) -> str:
    """
    Generate presentation outline using LLM
    
    Args:
        client: LLM client
        md_text: Markdown text
        image_messages: List of image messages
        model_name: Model name
        
    Returns:
        Generated presentation outline
    """
    prompt_template = load_prompt_template("GenerateSlidesOutlinePrompt.md")
    
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": prompt_template.format(content=md_text),
                },
            ] + image_messages,
        }
    ]
    

    response = client.chat.completions.create(
        model="gpt-4.1",
        messages=messages
    )
    return response.choices[0].message.content

def create_ppt_from_outline(outline: str, output_file: str, image_paths: List[str], pdf_base_name: str) -> None:
    """
    Create PowerPoint presentation from outline
    
    Args:
        outline: Presentation outline
        output_file: Output file name
        image_paths: List of image paths
        pdf_base_name: PDF file base name
    """
    prs = pptx.Presentation()
    slides = outline.split('---')

    for slide in slides:
        lines = slide.strip().split('\n')
        if not lines:
            continue

        slide_title = lines[0].replace('**', '').strip()
        slide_content = lines[1:]

        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_title

        content = slide.placeholders[1].text_frame
        for line in slide_content:
            line = line.strip()
            # Remove leading '- ' if present
            if line.startswith('- '):
                line = line[2:]
            if line:
                if '[Image' in line:
                    # Extract image number from the line
                    match = re.search(r'\[Image (\d+)\]', line)
                    if match:
                        image_number = int(match.group(1)) - 1
                        if image_number < len(image_paths):
                            image_path = str(Path('output') / pdf_base_name / 'auto' / image_paths[image_number])
                            try:
                                # Insert image below the text content
                                left = Inches(1)
                                top = Inches(2.5)
                                height = Inches(4)
                                slide.shapes.add_picture(image_path, left, top, height=height)
                            except Exception as e:
                                logger.error(f"Error adding image {image_path}: {e}")
                else:
                    p = content.add_paragraph()
                    p.text = line
                    p.level = 0

    prs.save(output_file)
    logger.info(f"PowerPoint presentation saved as '{output_file}'")

def display_usage_instructions():
    """Display usage instructions"""
    print("\n" + "="*80)
    print("PaperToSlides: Generate PowerPoint presentations from academic papers")
    print("="*80)
    print("\nUsage Examples:")
    print("  Basic usage:")
    print("    python GenerateSlidesOutline.py")
    print("\n  Specify a PDF file:")
    print("    python GenerateSlidesOutline.py --pdf path/to/your/paper.pdf")
    print("\n  Specify output filename:")
    print("    python GenerateSlidesOutline.py --output MyPresentation.pptx")
    print("\n  Combine options:")
    print("    python GenerateSlidesOutline.py --pdf path/to/paper.pdf --output MyPresentation.pptx")
    print("\nNote: Currently only the OpenAI model is implemented. To use other models,")
    print("you'll need to modify the code to include the appropriate API calls.")
    print("="*80)

def main():
    """Main function"""
    # Display usage instructions
    display_usage_instructions()

    # Parse command line arguments
    args = parse_arguments()
    
    # Setup API keys
    api_key = setup_api_keys(args)
    
    # Define file paths
    pdf_file_path = args.pdf
    local_md_dir = "output"
    os.makedirs(local_md_dir, exist_ok=True)
    
    # Get the expected markdown file path
    md_file_path = get_markdown_file_path(pdf_file_path, local_md_dir)
    pdf_base_name = Path(pdf_file_path).stem
    
    # Check if Markdown file already exists
    if os.path.exists(md_file_path):
        logger.info(f"Markdown file already exists: {md_file_path}, skipping PDF parsing")
    else:
        # If not, extract PDF content using function from MinerU.py
        logger.info(f"Markdown file does not exist, parsing PDF content...")
        md_file_path = extract_pdf_to_markdown(pdf_file_path, local_md_dir)
        # logger.error("PDF parsing functionality is disabled. Please uncomment the relevant code line or manually create the Markdown file.")
        # return
    
    # Load the extracted markdown content
    try:
        with open(md_file_path, "r", encoding="utf-8") as f:
            md_text = f.read()
    except Exception as e:
        logger.error(f"Error reading Markdown file {md_file_path}: {e}")
        return
    
    # Extract image paths from Markdown
    image_paths = extract_image_paths(md_text)
    
    # Prepare image messages
    image_messages = prepare_image_messages(image_paths, pdf_base_name)
    
    # Get LLM client, if you want to use Gemini, change to get_gemini_client(), set the model name to gemini-2.5-pro and add GEMINI_API_KEY to .env file
    client = get_openai_client(args.model, api_key)
    
    # Generate PowerPoint outline
    ppt_outline = generate_outline(client, md_text, image_messages, args.model)
    
    # Print formatted outline
    print(ppt_outline)
    
    # Create PowerPoint presentation
    create_ppt_from_outline(ppt_outline, args.output, image_paths, pdf_base_name)

if __name__ == "__main__":
    main()